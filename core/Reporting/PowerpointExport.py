from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.chart.data import CategoryChartData
import os
import core.Components.Utils as Utils
import copy
import six

class DummyProgressBar:
    def update(self):
        return

def replaceTextInParagraph(paragraph, search, replace):
    """
    Replace at a run level a text that will be searched in the given paragraph.
        Args:
            table: the table we want to replace a text inside
            search: The string to search for that will be replaced. //!\\ The search is done at a run level
            replace: The string to replace the search by.
        Returns:
            Return None if nothing was found, the run where the string was replaced otherwise
    """
    if search in paragraph.text:
        for run in paragraph.runs:
            inline = paragraph.runs
            # Loop added to work with runs (strings with same style)
            for i in range(len(inline)):
                if search in inline[i].text:
                    inline[i].text = inline[i].text.replace(search, replace)  
                    return run
        paragraph.text = paragraph.text.replace(search, replace)
        return paragraph
    return None


def findTextInTable(tbl, search):
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)
    for r in range(0, row_count):
        for c in range(0, col_count):
            cell = tbl.cell(r, c)
            paragraphs = cell.text_frame.paragraphs 
            for paragraph in paragraphs:
                if search in paragraph.text:
                    return paragraph
    return None

def replaceTextInTable(tbl, search, replace):
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)
    atLeastOneMatch = False
    for r in range(0, row_count):
        for c in range(0, col_count):
            cell = tbl.cell(r, c)
            paragraphs = cell.text_frame.paragraphs
            for paragraph in paragraphs:
                if search in paragraph.text:
                    res = replaceTextInParagraph(paragraph, search, replace)
                    if res is not None:
                        atLeastOneMatch = True
    return atLeastOneMatch

def replaceTextInChart(_chrt, _search, _replace):
    #chart_data = chrt.chart_data
    atLeastOneMatch = False
    return atLeastOneMatch

def findTextInTextFrame(shapeTextFrame, search):
    if (shapeTextFrame.text.find(search))!=-1:
        return shapeTextFrame.text.find(search)
    return None

def replaceTextInTextFrame(shapeTextFrame, search, replace):
    atLeastOneMatch = False
    if(shapeTextFrame.text.find(search))!=-1:
        text_frame = shapeTextFrame.text_frame
        paragraphs = text_frame.paragraphs
        for paragraph in paragraphs:
            if search in paragraph.text:
                res = replaceTextInParagraph(paragraph, search, replace)
                if res is not None:
                    atLeastOneMatch = True
    return atLeastOneMatch

def replaceTextInDocument(document, search, replace):
    """
    Replace at a run level a text that will be searched in all the document paragraphs and tables.
        Args:
            document: the document object (opened pptx with at lease one paragraph or table inside)
            search: The string to search for that will be replaced. //!\\ The search is done at a run level
            replace: The string to replace the search by.
    """
    for slide in document.slides:
        for shape in slide.shapes:
            if shape.has_table:
                replaceTextInTable(shape.table, search, replace)
            elif shape.has_text_frame:
                replaceTextInTextFrame(shape, search, replace)
            elif shape.has_chart:
                replaceTextInChart(shape.chart, search, replace)

def findTextInDocument(document, search):
    """
    Replace at a run level a text that will be searched in all the document paragraphs and tables.
        Args:
            document: the document object (opened pptx with at lease one paragraph or table inside)
            search: The string to search for that will be replaced. //!\\ The search is done at a run level
        Return:
            The slide index where the text was found, None otherwise
    """
    for slide_i, slide in enumerate(document.slides):
        for shape in slide.shapes:
            if shape.has_table:
                if findTextInTable(shape.table, search) is not None:
                    return slide_i
            elif shape.has_text_frame:
                if findTextInTextFrame(shape, search) is not None:
                    return slide_i
    return None

def findShapeContaining(document, slide_i, search):
    for shape in document.slides[slide_i].shapes:
        if shape.has_text_frame:
            if(shape.text.find(search))!=-1:
                return shape
    return None


def duplicate_slide(pres, index):
    # pylint: disable=protected-access
    """Duplicate the slide with the given index in pres.

    Adds slide to the end of the presentation"""
    source = pres.slides[index]
    # Tip: a nearly blank slide named TO_COPY
    # Keep placeholders for formatting but put a single white space inside
    dest = pres.slides.add_slide(SLD_LAYOUT_TO_COPY)
    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')
    try:
        for _, value in six.iteritems(source.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if not "notesSlide" in value.reltype:
                dest.rels.add_relationship(value.reltype, value._target, value.rId)
    except AttributeError:
        pass

def findSlideTableContaining(document, search):
    for i, slide in enumerate(document.slides):
        for shape in slide.shapes:
            if shape.has_table:
                if findTextInTable(shape.table, search) is not None:
                    return shape.table, i
    return None, None

def findTableInSlide(document, slide_i, search):
    for shape in document.slides[slide_i].shapes:
        if shape.has_table:
            if findTextInTable(shape.table, search) is not None:
                return shape.table

def copyShape(shape, idx):
    # pylint: disable=protected-access
    sp = shape._sp
    new_sp = None
    # ---duplicate original freeform---
    new_sp = copy.deepcopy(sp)
    new_shape = Shape(new_sp, shape._parent)
    # ---create a unique id for it---
    new_sp.nvSpPr.cNvPr.id = 1000 + idx
    # ---insert it after original---
    sp.addnext(new_sp)
    return new_shape

def deleteShape(shape):
    # pylint: disable=protected-access
    sp = shape._element
    sp.getparent().remove(sp)

def write_every_defect_fix(fixes, document, slide_i, count):
    # print("write fixes for "+str(count)+ " after "+str(last_defect_paragraph.text))
    shape_fix = None
    for fixe_i in range(len(fixes)):
        id_correctif = "A"+str(count)
        fixe = fixes[fixe_i]
        if len(fixes) > 1:
            id_correctif += "."+str(fixe_i+1)
        # Find the original fix to copy and fill
        shape_fix = findShapeContaining(document, slide_i, "var_c_id")
        if shape_fix is None:
            raise Exception("Error, fixe form text was deleted too soon. RIP")
        # Add notes paragraph after it
        h = shape_fix.height
        added_shape = copyShape(shape_fix, 1000+(count*100)+fixe_i)
        added_shape.top -= (h*(len(fixes)-fixe_i-1))
        replaceTextInTextFrame(added_shape, "var_c_id", id_correctif)
        replaceTextInTextFrame(added_shape, "var_c_title", fixe["title"])
    if shape_fix is not None:
        deleteShape(shape_fix)

def write_defect_from_input(result, document, table_d, slide_i, o_defect, count):
    desc = result[0]["description"].replace("\r", "")
    desc_paras = desc.split("\n")
    desc_paras += o_defect.get("notes", "").strip().split("\n")
    synthesis = result[0].get("synthesis", None)
    if synthesis is not None:
        desc_paras.insert(0, synthesis.replace("\r", ""))
    desc_paras.insert(1, "ToDo "+str(o_defect["redactor"]))
    replaceTextInTable(table_d, "var_d_description", desc_paras[0].strip())
    write_every_defect_fix(result[0]["fixes"], document, slide_i, count)
    return len(result[0]["fixes"])

def write_each_defect(document, defects_dict):
    """
    for each default
       Copy a table and a paragraph form the template marked with var_d_id and var_d_separator.
       replace the markers var_d_id var_d_separator var_d_id var_d_title var_d_ease var_d_impact var_d_description
    Then for each fixe of this default
        Copy a table form the template marked with var_c_id.
        replace the markers var_c_id var_d_separator var_c_title var_c_ease var_c_gain var_c_description

        Args:
            document: the document to search elements in
            defects_dict: the dictionary of defect gotten with the dedicated function getDefectDictFromExcel

    """
    levels = ["Critique", "Majeur", "Important", "Mineur"]
    count_defects = 0
    count_fixes = 0
    total_len = 0
    for level in levels:
        total_len += len(defects_dict[level].values())
    for level in levels:
        level_count = 0
        _, slide_copy_i = findSlideTableContaining(document, "var_d_id")
        # COPY SLIDE X TIME, SEPARATED
        for i in range(1, len(defects_dict[level].values())):
            duplicate_slide(document, slide_copy_i)
            new_slide_i = slide_copy_i+i
            move_slide(document, -1, new_slide_i)
        # FILL COPIED SLIDES
        for defect_dict in defects_dict[level].values():
            new_slide_i = slide_copy_i + level_count
            level_count += 1
            count_defects += 1
            o_defect = defect_dict["description"]
            table_d = findTableInSlide(document, new_slide_i, "var_d_id")
            replaceTextInTable(table_d, "var_d_id", "D"+str(count_defects))
            replaceTextInTable(table_d, "var_d_title", o_defect["title"])
            replaceTextInTable(table_d, "var_d_ease", o_defect["ease"])
            replaceTextInTable(table_d, "var_d_impact", o_defect["impact"])
            
            result = [
                    {
                        "id":"0",
                        "title": o_defect["title"],
                        "ease": o_defect["ease"],
                        "impact": o_defect["impact"],
                        "risk": o_defect["risk"],
                        "type": o_defect["type"],
                        "description": "Description",
                        "notes": o_defect.get("notes", ""),
                        "fixes": [
                            {
                                "title": "A def",
                                "execution": "Modérée",
                                "gain": "Moyen",
                                "description": "A definir"
                            }
                        ]
                    }
                ]
            print("Writing defect "+str(count_defects)+"")
            count_fixes += write_defect_from_input(result, document, table_d, new_slide_i, o_defect, count_defects)
            progressbar.update()
        if level_count == 0:
            table_d = findTableInSlide(document, slide_copy_i, "var_d_id")
            if table_d is not None:
                replaceTextInTable(table_d, "var_d_id", "TO_DELETE")
    _, slide_copy_i = findSlideTableContaining(document, "TO_DELETE")
    while slide_copy_i is not None:
        delete_slide(document, slide_copy_i)
        _, slide_copy_i = findSlideTableContaining(document, "TO_DELETE")
    return count_defects, count_fixes


def move_slide(presentation, old_index, new_index):
    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def delete_slide(presentation, index):
    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    xml_slides.remove(slides[index])

def addSerieToChart(presentation, index_chart, serie_name, serie):
    count_chart = 0
    for _slide_i, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if shape.has_chart:
                if count_chart == index_chart:
                    chart_data = CategoryChartData()
                    chart_data.categories = ("Critique", "Majeur", "Important", "Mineur")
                    chart_data.add_series(serie_name, serie)
                    shape.chart.replace_data(chart_data)
                    shape.chart.value_axis.maximum_scale = max(serie)+1
                    return True
                count_chart += 1
    return False

def createReport(defects_dict, template, out_name, **kwargs):
    document = Presentation(template)
    global progressbar
    global SLD_LAYOUT_TO_COPY
    SLD_LAYOUT_TO_COPY = document.slide_layouts.get_by_name("TO_COPY")
    if SLD_LAYOUT_TO_COPY is None:
        raise Exception("The pptx template does not contain a TO_COPY layout")
    client_name = kwargs.get("client", "").strip()
    progressbar = kwargs.get("progressbar", DummyProgressBar())
    total_len = 0
    levels = ["Critique", "Majeur", "Important", "Mineur"]
    for level in levels:
        total_len += len(defects_dict[level].values())
    nb_steps = total_len # 1 step by defect
    nb_steps += 1 # step for general stuff
    nb_steps += 1 # step for saving
    progressbar.show(nb_steps)
    if client_name != "":
        replaceTextInDocument(document, "var_client", client_name)
    contract_name = kwargs.get("contract", "").strip()
    if contract_name != "":
        if contract_name.lower().startswith("test "):
            contract_name = "du t"+contract_name[1:]
        elif contract_name.lower().startswith("audit "):
            contract_name = "de l'a"+contract_name[1:]
        replaceTextInDocument(document, "var_contract", contract_name)
    nb_critical = len(defects_dict["Critique"].keys())
    nb_major = len(defects_dict["Majeur"].keys())
    nb_important = len(defects_dict["Important"].keys())
    nb_minor = len(defects_dict["Mineur"].keys())
    replaceTextInDocument(document, "var_nb_d_critical", str(nb_critical))
    replaceTextInDocument(document, "var_nb_d_major", str(nb_major))
    replaceTextInDocument(document, "var_nb_d_important", str(nb_important))
    replaceTextInDocument(document, "var_nb_d_minor", str(nb_minor))
    addSerieToChart(document, 0, 'Criticity', (nb_critical,nb_major,nb_important,nb_minor))
    progressbar.update()
    print("Write each defect ...")
    count_defect, count_fixes = write_each_defect(document, defects_dict)
    replaceTextInDocument(document, "var_nb_d_total", str(count_defect))
    replaceTextInDocument(document, "var_nb_fix", str(count_fixes))

    #print("Write services table ...")
    #populate_services_table(document, kwargs.get("root", None))
    #print("Converting Markdown ...")
    #markdownToWordInDocument(document)
    #print("Saving ...")
    # After that, we add the previously copied table
    #paragraph._p.addnext(new_tbl)
    dir_path = os.path.dirname(os.path.realpath(__file__))
    out_path = os.path.join(dir_path, "../../exports/", out_name+".pptx")
    document.save(out_path)
    progressbar.update()
    print("Generated report at "+str(out_path))
